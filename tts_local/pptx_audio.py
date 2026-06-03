"""PPTX 슬라이드에 오디오를 삽입하고 자동재생(onBegin) 타이밍을 설정한다.
python-pptx의 add_movie + 직접 XML 조작 방식 (TTS_google 검증 로직 기반)."""
from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

_P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def _mime(path: Path) -> str:
    s = str(path).lower()
    if s.endswith(".mp3"):
        return "audio/mpeg"
    if s.endswith(".m4a"):
        return "audio/mp4"
    return "audio/wav"


def _set_autoplay(slide, audio_shape) -> None:
    """오디오를 슬라이드 진입 시 자동재생되도록 timing 트리 추가."""
    try:
        children = list(audio_shape.element)
        if children:
            children[0].set("action", "play")
    except Exception as e:  # noqa: BLE001
        logger.debug("autoplay action 설정 실패: %s", e)

    try:
        from lxml import etree
        timing = etree.SubElement(slide.element, f"{_P_NS}timing")
        tn_lst = etree.SubElement(timing, f"{_P_NS}tnLst")
        par = etree.SubElement(tn_lst, f"{_P_NS}par")
        ctn = etree.SubElement(par, f"{_P_NS}cTn")
        ctn.set("id", "1"); ctn.set("dur", "indefinite"); ctn.set("restart", "never")
        child_tn_lst = etree.SubElement(ctn, f"{_P_NS}childTnLst")
        audio_el = etree.SubElement(child_tn_lst, f"{_P_NS}audio")
        c_media = etree.SubElement(audio_el, f"{_P_NS}cMediaNode")
        c_tn = etree.SubElement(c_media, f"{_P_NS}cTn")
        c_tn.set("id", "2"); c_tn.set("fill", "hold"); c_tn.set("display", "0")
        st_cond = etree.SubElement(c_tn, f"{_P_NS}stCondLst")
        cond = etree.SubElement(st_cond, f"{_P_NS}cond")
        cond.set("delay", "0"); cond.set("evt", "onBegin")
        tgt = etree.SubElement(c_media, f"{_P_NS}tgt")
        sp_tgt = etree.SubElement(tgt, f"{_P_NS}spTgt")
        sp_tgt.set("spid", str(audio_shape.shape_id))
    except Exception as e:  # noqa: BLE001
        logger.debug("autoplay timing 설정 실패: %s", e)


def embed_audio(pptx_path: Path, audio_by_slide: dict[int, Path], out_path: Path,
                autoplay: bool = True) -> Path:
    """audio_by_slide: {1-based 덱 슬라이드번호: 오디오 경로}. 결과 PPTX 저장."""
    prs = Presentation(str(pptx_path))
    # 작은 스피커 아이콘을 슬라이드 밖(왼쪽)으로 배치해 화면을 가리지 않게 함
    left, top, w, h = Inches(-1.0), Inches(0.0), Inches(0.3), Inches(0.3)
    inserted = 0
    for idx, slide in enumerate(prs.slides, start=1):
        audio = audio_by_slide.get(idx)
        if not audio or not Path(audio).exists():
            continue
        shape = slide.shapes.add_movie(str(audio), left, top, w, h, mime_type=_mime(Path(audio)))
        if autoplay:
            try:
                shape.media_format.play_settings.play_on_click = False
                shape.media_format.play_settings.play_automatically = True
            except Exception:  # noqa: BLE001
                pass
            _set_autoplay(slide, shape)
        inserted += 1
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    logger.info("오디오 %d개 삽입 → %s", inserted, out_path)
    return out_path
